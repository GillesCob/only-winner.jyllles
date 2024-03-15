from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from datetime import datetime

#ATTENTION : BIEN LIRE LES INFOS SUIVANTES : 
#Va créer les cartes dans le ppt
#Créer le dossier du mois et y mettre le ppt avec le bon nombre de slides
#Mettre à jour le numéro de la feuille Excel utilisée pour récupérer les datas
#Mettre à jour le mois concerné !
#Mettre à jour le nombre_events

actual_year = str(datetime.now().year)
actual_month = str(datetime.now().month)
actual_day = str(datetime.now().day)
scrapping_month = "Mars"
Feuille_datas = str(actual_day)
diapositive_xx = 1

valeurs_colonne_J = []
comp_of_sport_list = []
comp_of_sport_card_list = []
valeurs_colonne_L = []
# Charger le fichier Excel
excel_month_file_path = f'/Users/gillescobigo/Documents/Gilles/Dev/Only Winners/DATAS/2024/{scrapping_month}/EXCEL/DATAS 2.xlsx'
excel_BDD_INITIALE_file_path = f'/Users/gillescobigo/Documents/Gilles/Dev/Only Winners/DATAS/BDD INITIALE.xlsx'

excel_month_workbook = load_workbook(excel_month_file_path)
excel_BDD_INITIALE_workbook = load_workbook(excel_BDD_INITIALE_file_path)
excel_sheet = excel_month_workbook[Feuille_datas]
nom_comp_of_sport_sheet = excel_BDD_INITIALE_workbook["COMP OF SPORT"]


# Charger la présentation PowerPoint
pptx_file_path = f"/Users/gillescobigo/Documents/Gilles/Dev/Only Winners/DATAS/2024/{scrapping_month}/POWERPOINT/cartes_sans_images.pptx"
presentation = Presentation(pptx_file_path)

NOMBRE_EVENTS = diapositive_xx + 2 #Mettre 2 de plus que le nombre d'events

#Je récupère toutes les valeurs présentes dans la colonne A
for cell in nom_comp_of_sport_sheet['A'][1:]:
    # Vérifier si la cellule n'est pas vide et ajouter sa valeur à la liste
    if cell.value is not None:
        comp_of_sport_list.append(cell.value)

#Je récupère toutes les valeurs présentes dans la colonne C
for cell in nom_comp_of_sport_sheet['C'][1:]:
    # Vérifier si la cellule n'est pas vide et ajouter sa valeur à la liste
    if cell.value is not None:
        comp_of_sport_card_list.append(cell.value)

# Récupérer la valeur de la cellule A2 dans Excel
for i in range (2, NOMBRE_EVENTS):

    number = excel_sheet[f'A{i}'].value
    competition_sport = f'{excel_sheet[f"F{i}"].value} of {excel_sheet[f"E{i}"].value}'
    prompt = excel_sheet[f'L{i}'].value
    winner = excel_sheet[f'I{i}'].value
    ref = excel_sheet[f'A{i}'].value

    # Accéder aux diapos en fonction du i
    slides = presentation.slides[i-2]

    # # Trouver la zone de texte avec le texte "Date"
    for slide in slides.shapes:

        #Compétition de sport
        if slide.has_text_frame and "Compétition de sport" in slide.text_frame.text:
            if competition_sport is not None :
                if competition_sport in comp_of_sport_list:
                    j = comp_of_sport_list.index(competition_sport)
                    new_competition_sport = comp_of_sport_card_list[j]
                    slide.text_frame.text = new_competition_sport
                    
                else:
                    slide.text_frame.text = competition_sport
                    print(f"PAS DE TRADUCTION : {competition_sport}")

        #Prompt
        elif slide.has_text_frame and "Prompt" in slide.text_frame.text:
            # Mettre à jour le texte avec la valeur de l'Excel
            if prompt is not None :
                slide.text_frame.text = prompt
                
        #Ref
        elif slide.has_text_frame and "REF" in slide.text_frame.text:
            # Mettre à jour le texte avec la valeur de l'Excel
            if ref is not None :
                slide.text_frame.text = ref
                
        #Note
        notes_slide = slides.notes_slide
        notes_placeholder = notes_slide.notes_text_frame
        if number is not None:
            notes_placeholder.text = f'{actual_year}-{actual_month}-{number}'

        # Définir la taille de la police à 50 points
        if hasattr(slide, "text_frame"):
            for paragraph in slide.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER  # Centrage
                for run in paragraph.runs:
                    run.font.size = Pt(50)
                    
    # Lien vers l'image à insérer sur sa diapo
    image_path = f'/Users/gillescobigo/Documents/Gilles/Dev/Only Winners/DATAS/2024/{scrapping_month}/IMAGES_MIDJOURNEY/{prompt}.png'
    #print(winner)

    #presentation.slides[0].shapes[-1].element.clear()
    try:
        presentation.slides[i-2].shapes.add_picture(image_path, Inches(0.78), Inches(2.8), width=Inches(14.19))
    except FileNotFoundError:
        print(f"{prompt} ,oil painting style, colorful lighting,--ar 1:1 --c 100 --s 965 --v 5.1")
        
print(f"Création du PPT terminée !")   

# Enregistrer la présentation PowerPoint mise à jour
presentation.save(f"/Users/gillescobigo/Documents/Gilles/Dev/Only Winners/DATAS/2024/{scrapping_month}/POWERPOINT/cartes_avec_images.pptx")

# Fermer le fichier Excel
excel_month_workbook.close()
excel_BDD_INITIALE_workbook.close()