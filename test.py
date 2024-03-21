from pptx import Presentation
from PIL import Image

def export_slides_to_png(pptx_file, output_folder):
    prs = Presentation(pptx_file)

    for i, slide in enumerate(prs.slides):
        image = Image.new("RGB", (prs.slide_width, prs.slide_height), "white")
        image_path = f"{output_folder}/slide_{i+1}.png"

        for shape in slide.shapes:
            if hasattr(shape, "image"):
                with open(f"{output_folder}/image{i}.png", "wb") as f:
                    f.write(shape.image.blob)

        image.save(image_path)

pptx_file = "/Users/gillescobigo/Desktop/cartes_avec_images.pptx"
output_folder = "/Users/gillescobigo/Desktop"

export_slides_to_png(pptx_file, output_folder)
